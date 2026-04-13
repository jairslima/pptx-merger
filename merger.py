"""
PPTX Merger by Jair Lima
Módulo responsável pela mesclagem de arquivos PPTX.
"""

import copy
import itertools
import os
from pptx import Presentation
from pptx.opc.package import Part
from pptx.opc.package import PackURI
from lxml import etree

# Contador global para nomes únicos de mídia no pacote mesclado
_media_counter = itertools.count(1)

# Namespaces XML do PPTX
NS_P   = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_DGM = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'

# Atributos que contêm IDs de relacionamento
REL_ATTRS = (
    f'{{{NS_R}}}embed',
    f'{{{NS_R}}}link',
    f'{{{NS_R}}}id',
    f'{{{NS_R}}}href',
)

# Tipos de relacionamento a ignorar (pertencentes à estrutura do slide, não ao conteúdo)
SKIP_REL_TYPES = (
    'slideLayout',
    'notesSlide',
    'slide/',
)


def get_slide_count(filepath: str) -> int:
    """Retorna o número de slides de um arquivo PPTX."""
    try:
        prs = Presentation(filepath)
        return len(prs.slides)
    except Exception:
        return 0


def merge_pptx(file_list: list, output_path: str,
               progress_callback=None) -> int:
    """
    Mescla uma lista de arquivos PPTX em um único arquivo.

    Args:
        file_list:         Lista ordenada de caminhos para os arquivos PPTX.
        output_path:       Caminho do arquivo de saída.
        progress_callback: Função opcional (atual, total) chamada a cada slide copiado.

    Returns:
        Número total de slides no arquivo mesclado.
    """
    if not file_list:
        raise ValueError("Nenhum arquivo fornecido para mesclagem.")

    # Contar total de slides (excluindo o primeiro arquivo, que é a base)
    total_extra = sum(
        get_slide_count(f) for f in file_list[1:]
    )

    base_prs = Presentation(file_list[0])
    copied = 0

    for filepath in file_list[1:]:
        src_prs = Presentation(filepath)
        for slide in src_prs.slides:
            _copy_slide(base_prs, slide)
            copied += 1
            if progress_callback:
                progress_callback(copied, total_extra)

    _ensure_output_dir(output_path)
    base_prs.save(output_path)
    return len(base_prs.slides)


def _get_blank_layout(prs: Presentation):
    """Retorna o layout 'em branco' da apresentação, com fallback seguro."""
    blank_names = {'blank', 'em branco', 'vazio', 'branco', 'leer', 'vuoto'}
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() in blank_names:
            return layout
    # fallback: último layout (geralmente é o mais vazio)
    return prs.slide_layouts[-1]


def _clone_part(src_part) -> Part:
    """
    Clona uma parte de mídia (imagem, vídeo, áudio) atribuindo um nome
    interno único, evitando colisão de nomes ao mesclar apresentações
    que possuem arquivos de mídia com nomes idênticos.
    """
    ext = os.path.splitext(str(src_part.partname))[1] or '.bin'
    idx = next(_media_counter)
    new_name = PackURI(f'/ppt/media/merged_{idx}{ext}')
    return Part(new_name, src_part.content_type, src_part.blob)


def _copy_slide(dest_prs: Presentation, src_slide) -> None:
    """
    Copia um slide (com imagens, formas e plano de fundo) para dest_prs.
    """
    blank_layout = _get_blank_layout(dest_prs)
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # Construir mapa: rId_antigo -> rId_novo para todas as relações de conteúdo
    rId_map = {}
    for old_rId, rel in src_slide.part.rels.items():
        # Pular relações estruturais (layout, notas, etc.)
        if any(skip in rel.reltype for skip in SKIP_REL_TYPES):
            continue

        try:
            if rel.is_external:
                new_rId = new_slide.part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                # Clonar a parte de mídia com nome único para evitar
                # colisão entre arquivos que têm imagens de mesmo nome interno
                cloned = _clone_part(rel.target_part)
                new_rId = new_slide.part.relate_to(cloned, rel.reltype)
            rId_map[old_rId] = new_rId
        except Exception:
            pass  # relação incompatível ou duplicada; ignorar

    # Deep copy do elemento raiz do slide (<p:sld>)
    src_el = copy.deepcopy(src_slide._element)

    # Remapear todos os atributos de relacionamento na cópia
    for el in src_el.iter():
        for attr in REL_ATTRS:
            if attr in el.attrib and el.attrib[attr] in rId_map:
                el.attrib[attr] = rId_map[el.attrib[attr]]

    # ---- Copiar spTree (formas e conteúdo) ----
    src_cSld = src_el.find(f'{{{NS_P}}}cSld')
    if src_cSld is not None:
        src_sp_tree = src_cSld.find(f'{{{NS_P}}}spTree')
        if src_sp_tree is not None:
            dest_sp_tree = new_slide.shapes._spTree
            # Remover placeholders do layout em branco
            for child in list(dest_sp_tree):
                dest_sp_tree.remove(child)
            # Inserir conteúdo do slide de origem
            for child in src_sp_tree:
                dest_sp_tree.append(copy.deepcopy(child))

    # ---- Copiar plano de fundo (se definido explicitamente no slide) ----
    src_bg = src_cSld.find(f'{{{NS_P}}}bg') if src_cSld is not None else None
    if src_bg is not None:
        dest_cSld = new_slide._element.find(f'{{{NS_P}}}cSld')
        if dest_cSld is not None:
            existing_bg = dest_cSld.find(f'{{{NS_P}}}bg')
            if existing_bg is not None:
                dest_cSld.remove(existing_bg)
            dest_cSld.insert(0, copy.deepcopy(src_bg))

    # ---- Copiar transições de slide ----
    src_transition = src_el.find(f'{{{NS_P}}}transition')
    if src_transition is not None:
        existing = new_slide._element.find(f'{{{NS_P}}}transition')
        if existing is not None:
            new_slide._element.remove(existing)
        new_slide._element.append(copy.deepcopy(src_transition))


def _ensure_output_dir(path: str) -> None:
    """Garante que o diretório de saída existe."""
    directory = os.path.dirname(path)
    if directory and not os.path.exists(directory):
        os.makedirs(directory, exist_ok=True)


def unique_output_path(path: str) -> str:
    """
    Se o arquivo já existir, adiciona sufixo _1, _2... até encontrar nome livre.
    """
    if not os.path.exists(path):
        return path
    base, ext = os.path.splitext(path)
    counter = 1
    while True:
        candidate = f"{base}_{counter}{ext}"
        if not os.path.exists(candidate):
            return candidate
        counter += 1
